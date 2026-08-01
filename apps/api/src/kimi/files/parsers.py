"""Document parsers.

Each parser returns segments that know where they came from. The audited
defects in ``legacy_streamlit/core/file_tools.py`` are fixed here and every one
has a regression test:

* **Scanned PDFs vanished.** Pages with no text layer produced ``""``, which the
  context builder then dropped entirely, so the model never learned the file
  existed. Now detected explicitly and reported as ``NO_TEXT_LAYER``.
* **Encrypted PDFs crashed.** ``reader.decrypt`` was never called.
* **PPTX lost tables, groups and notes.** ``hasattr(shape, "text")`` is false for
  a ``GraphicFrame``, which is what a table is, so table-driven decks — the
  normal case for financial slides — came back nearly empty.
* **DOCX destroyed reading order.** All paragraphs were emitted, then all
  tables, so a table on page 2 landed after 40 pages of prose.
* **CSV lied about size.** ``nrows=500`` was reported as "Rows loaded: 500" with
  no indication the file had two million.
* **Eight unmarked truncations.** Text was hard-sliced with no sentinel, so a
  cut mid-page looked like the end of the page.
"""

from __future__ import annotations

import base64
import csv
import io
import re
from collections.abc import Iterator
from typing import Any, Final

from kimi.files.models import (
    DocumentKind,
    ParsedDocument,
    ParseStatus,
    RefKind,
    Segment,
    SegmentRef,
)

#: Per-segment cap. Exceeding it appends an explicit sentinel.
MAX_SEGMENT_CHARS: Final = 6_000
#: Per-document cap across all segments.
MAX_DOCUMENT_CHARS: Final = 60_000
TRUNCATION_SENTINEL: Final = "\n[… truncated]"

#: A page with fewer than this many characters is treated as having no text
#: layer. Real prose pages clear it comfortably; scanned pages return 0-2.
MIN_PAGE_TEXT_CHARS: Final = 12

MAX_PDF_PAGES: Final = 300
MAX_SHEET_ROWS: Final = 500
MAX_CSV_SAMPLE_ROWS: Final = 200

_WS = re.compile(r"[ \t\r\f\v]+")
_BLANKS = re.compile(r"\n{3,}")


class ParseFailure(Exception):
    """Raised when a file cannot be parsed. The message is user-safe."""


def decode_text(raw: bytes) -> str:
    """The prototype's encoding ladder, kept — it handles Arabic CP1256 well."""
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "cp1256", "cp1252", "latin-1"):
        try:
            return raw.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="replace")


def tidy(text: str) -> str:
    text = _WS.sub(" ", text or "")
    text = _BLANKS.sub("\n\n", text)
    return "\n".join(line.strip() for line in text.splitlines()).strip()


def clip(text: str, limit: int = MAX_SEGMENT_CHARS) -> tuple[str, bool]:
    """Cut to ``limit`` and say so. Never a silent slice."""
    if len(text) <= limit:
        return text, False
    cut = text[:limit]
    # Prefer a sentence or word boundary so the cut does not land mid-word.
    for boundary in ("\n", ". ", " "):
        idx = cut.rfind(boundary)
        if idx > limit * 0.6:
            cut = cut[: idx + len(boundary)]
            break
    return cut.rstrip() + TRUNCATION_SENTINEL, True


def _budgeted(segments: list[Segment]) -> tuple[list[Segment], bool]:
    """Apply the document-wide cap, charging the full segment cost."""
    kept: list[Segment] = []
    used = 0
    truncated = False
    for segment in segments:
        # Charge the label too — the prototype charged only the body, so the
        # nominal ceiling was exceeded by an unbounded margin.
        cost = len(segment.text) + len(segment.ref.label) + 4
        if used + cost > MAX_DOCUMENT_CHARS:
            remaining = MAX_DOCUMENT_CHARS - used - len(segment.ref.label) - 4
            if remaining > 200:
                text, _ = clip(segment.text, remaining)
                kept.append(Segment(ref=segment.ref, text=text, truncated=True))
            truncated = True
            break
        kept.append(segment)
        used += cost
    return kept, truncated


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def parse_pdf(doc: ParsedDocument, data: bytes) -> ParsedDocument:
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError

    try:
        reader = PdfReader(io.BytesIO(data))
    except (PdfReadError, OSError, ValueError) as exc:
        raise ParseFailure("This PDF could not be opened. It may be damaged.") from exc

    if reader.is_encrypted:
        # The prototype was not encryption-aware at all and simply produced
        # empty text for every page.
        try:
            if reader.decrypt("") == 0:
                doc.status = ParseStatus.PASSWORD_REQUIRED
                doc.summary = f"{doc.filename}: this PDF is password-protected and was not read."
                return doc
        except (PdfReadError, NotImplementedError) as exc:
            doc.status = ParseStatus.ENCRYPTED
            doc.summary = f"{doc.filename}: this PDF uses encryption this app cannot open."
            doc.warnings.append(str(exc)[:120])
            return doc

    total_pages = len(reader.pages)
    doc.metadata["pages"] = total_pages
    if total_pages == 0:
        doc.status = ParseStatus.FAILED
        doc.summary = f"{doc.filename}: the PDF contains no pages."
        return doc

    segments: list[Segment] = []
    pages_with_text = 0
    read_pages = min(total_pages, MAX_PDF_PAGES)

    for index in range(read_pages):
        try:
            raw = reader.pages[index].extract_text() or ""
        except Exception:
            doc.warnings.append(f"Page {index + 1} could not be read.")
            continue
        text = tidy(raw)
        if len(text) >= MIN_PAGE_TEXT_CHARS:
            pages_with_text += 1
            body, was_cut = clip(text)
            segments.append(
                Segment(ref=SegmentRef(RefKind.PAGE, index + 1), text=body, truncated=was_cut)
            )

    if read_pages < total_pages:
        doc.warnings.append(f"Only the first {read_pages} of {total_pages} pages were read.")

    # The load-bearing fix: pages exist but none carry text.
    if pages_with_text == 0:
        doc.status = ParseStatus.NO_TEXT_LAYER
        doc.summary = (
            f"{doc.filename}: {total_pages} page(s), but no selectable text. "
            "This looks like a scanned document. OCR is not available in this "
            "build, so its contents could not be read."
        )
        return doc

    doc.segments, over_budget = _budgeted(segments)
    doc.metadata["pages_with_text"] = pages_with_text
    doc.status = (
        ParseStatus.PARTIAL
        if over_budget or pages_with_text < total_pages or doc.warnings
        else ParseStatus.PARSED
    )
    doc.summary = f"{doc.filename}: PDF, {total_pages} page(s), text found on {pages_with_text}."
    return doc


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def _docx_blocks(document: Any) -> Iterator[tuple[str, Any]]:
    """Yield paragraphs and tables in true document order.

    python-docx exposes ``.paragraphs`` and ``.tables`` as separate lists, which
    is how the prototype ended up emitting all prose then all tables. Walking the
    body element preserves the order the author wrote.
    """
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = document.element.body
    for child in body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            yield "p", Paragraph(child, document)
        elif tag == "tbl":
            yield "tbl", Table(child, document)


def _table_to_text(table: Any) -> str:
    rows = []
    for row in table.rows:
        cells = [tidy(c.text) for c in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def parse_docx(doc: ParsedDocument, data: bytes) -> ParsedDocument:
    from docx import Document

    try:
        document = Document(io.BytesIO(data))
    except Exception as exc:
        raise ParseFailure("This Word document could not be opened.") from exc

    segments: list[Segment] = []
    buffer: list[str] = []
    section = 1
    tables = 0

    def flush() -> None:
        nonlocal section
        text = tidy("\n".join(buffer))
        buffer.clear()
        if text:
            body, was_cut = clip(text)
            segments.append(
                Segment(ref=SegmentRef(RefKind.SECTION, section), text=body, truncated=was_cut)
            )
            section += 1

    for kind, block in _docx_blocks(document):
        if kind == "p":
            text = block.text.strip()
            if not text:
                continue
            style = (getattr(block.style, "name", "") or "").lower()
            if style.startswith("heading"):
                # A heading starts a new citable section.
                flush()
                buffer.append(f"## {text}")
            else:
                buffer.append(text)
        else:
            flush()
            tables += 1
            body, was_cut = clip(_table_to_text(block))
            if body:
                segments.append(
                    Segment(ref=SegmentRef(RefKind.TABLE, tables), text=body, truncated=was_cut)
                )
    flush()

    if not segments:
        doc.status = ParseStatus.NO_TEXT_LAYER
        doc.summary = f"{doc.filename}: the document contains no readable text."
        return doc

    doc.segments, over_budget = _budgeted(segments)
    doc.metadata |= {"sections": section - 1, "tables": tables}
    doc.status = ParseStatus.PARTIAL if over_budget else ParseStatus.PARSED
    doc.summary = f"{doc.filename}: Word document, {section - 1} section(s), {tables} table(s)."
    return doc


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def _shape_texts(shape: Any) -> Iterator[str]:
    """Recurse through groups and read tables.

    ``hasattr(shape, "text")`` is False for a GraphicFrame, which is what a
    table is, and groups were never entered — so the prototype returned almost
    nothing for table-driven decks.
    """
    shape_type = getattr(shape, "shape_type", None)
    # 6 == GROUP in MSO_SHAPE_TYPE; compare by name to avoid the enum import.
    if str(shape_type).startswith("GROUP") or hasattr(shape, "shapes"):
        for child in getattr(shape, "shapes", []):
            yield from _shape_texts(child)
        return

    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            cells = [tidy(c.text) for c in row.cells]
            if any(cells):
                yield " | ".join(cells)
        return

    if getattr(shape, "has_text_frame", False):
        text = tidy(shape.text_frame.text)
        if text:
            yield text
        return

    text = tidy(getattr(shape, "text", "") or "")
    if text:
        yield text


def parse_pptx(doc: ParsedDocument, data: bytes) -> ParsedDocument:
    from pptx import Presentation

    try:
        deck = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ParseFailure("This PowerPoint file could not be opened.") from exc

    segments: list[Segment] = []
    slides = 0
    notes_found = 0

    for index, slide in enumerate(deck.slides, start=1):
        slides = index
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_texts(shape))
        text = tidy("\n".join(parts))
        if text:
            body, was_cut = clip(text)
            segments.append(
                Segment(ref=SegmentRef(RefKind.SLIDE, index), text=body, truncated=was_cut)
            )

        # Speaker notes are often the substantive argument and were never read.
        if slide.has_notes_slide:
            notes = tidy(slide.notes_slide.notes_text_frame.text)
            if notes:
                notes_found += 1
                body, was_cut = clip(notes)
                segments.append(
                    Segment(ref=SegmentRef(RefKind.NOTES, index), text=body, truncated=was_cut)
                )

    if not segments:
        doc.status = ParseStatus.NO_TEXT_LAYER
        doc.summary = f"{doc.filename}: {slides} slide(s), none containing readable text."
        return doc

    doc.segments, over_budget = _budgeted(segments)
    doc.metadata |= {"slides": slides, "slides_with_notes": notes_found}
    doc.status = ParseStatus.PARTIAL if over_budget else ParseStatus.PARSED
    doc.summary = (
        f"{doc.filename}: presentation, {slides} slide(s), {notes_found} with speaker notes."
    )
    return doc


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def parse_xlsx(doc: ParsedDocument, data: bytes) -> ParsedDocument:
    from openpyxl import load_workbook

    try:
        book = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        raise ParseFailure("This spreadsheet could not be opened.") from exc

    segments: list[Segment] = []
    sheet_info: list[dict[str, Any]] = []

    try:
        for sheet in book.worksheets:
            total_rows = sheet.max_row or 0
            total_cols = sheet.max_column or 0
            lines: list[str] = []
            read = 0
            for row in sheet.iter_rows(max_row=min(total_rows, MAX_SHEET_ROWS), values_only=True):
                cells = ["" if v is None else str(v).strip() for v in row]
                if any(cells):
                    lines.append(" | ".join(cells))
                read += 1

            sheet_info.append(
                {"name": sheet.title, "rows": total_rows, "columns": total_cols, "read": read}
            )
            if not lines:
                continue

            header = f"Rows 1-{read} of {total_rows}"
            if total_rows > read:
                header += " (sampled)"
            body, was_cut = clip(f"{header}\n" + "\n".join(lines))
            segments.append(
                Segment(
                    ref=SegmentRef(RefKind.SHEET, 0, sheet.title),
                    text=body,
                    truncated=was_cut or total_rows > read,
                )
            )
            if total_rows > read:
                doc.warnings.append(f"Sheet {sheet.title}: {total_rows} rows, first {read} read.")
    finally:
        book.close()

    doc.metadata["sheets"] = sheet_info
    if not segments:
        doc.status = ParseStatus.NO_TEXT_LAYER
        doc.summary = f"{doc.filename}: the workbook contains no readable cells."
        return doc

    doc.segments, over_budget = _budgeted(segments)
    doc.status = ParseStatus.PARTIAL if over_budget or doc.warnings else ParseStatus.PARSED
    names = ", ".join(s["name"] for s in sheet_info)
    doc.summary = f"{doc.filename}: spreadsheet with {len(sheet_info)} sheet(s): {names}."
    return doc


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------


def parse_csv(doc: ParsedDocument, data: bytes) -> ParsedDocument:
    text = decode_text(data)
    sample = text[:8192]
    try:
        dialect: Any = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        # The prototype had no delimiter sniffing, so a semicolon CSV parsed
        # into a single column.
        dialect = csv.excel

    reader = csv.reader(io.StringIO(text), dialect)
    try:
        rows = list(reader)
    except csv.Error as exc:
        raise ParseFailure("This CSV could not be parsed.") from exc

    if not rows:
        doc.status = ParseStatus.NO_TEXT_LAYER
        doc.summary = f"{doc.filename}: the file is empty."
        return doc

    header = rows[0]
    body_rows = rows[1:]
    # The true count, not the sample size. The prototype reported the sample as
    # the total, so the model confidently stated wrong dataset sizes.
    total_rows = len(body_rows)
    shown = body_rows[:MAX_CSV_SAMPLE_ROWS]

    lines = [" | ".join(header)]
    lines += [" | ".join(r) for r in shown]
    note = f"Total data rows: {total_rows}. Showing {len(shown)}."
    body, was_cut = clip(f"{note}\n" + "\n".join(lines))

    doc.segments = [
        Segment(
            ref=SegmentRef(RefKind.ROW, 0, f"rows 1-{len(shown)}"),
            text=body,
            truncated=was_cut or total_rows > len(shown),
        )
    ]
    doc.metadata |= {
        "columns": header,
        "column_count": len(header),
        "total_rows": total_rows,
        "sampled_rows": len(shown),
        "delimiter": getattr(dialect, "delimiter", ","),
    }
    if total_rows > len(shown):
        doc.warnings.append(f"{total_rows} rows total; the first {len(shown)} were included.")
    doc.status = ParseStatus.PARTIAL if total_rows > len(shown) else ParseStatus.PARSED
    doc.summary = f"{doc.filename}: CSV with {len(header)} column(s) and {total_rows} data row(s)."
    return doc


# ---------------------------------------------------------------------------
# Plain text
# ---------------------------------------------------------------------------


def parse_text(doc: ParsedDocument, data: bytes) -> ParsedDocument:
    text = tidy(decode_text(data))
    if not text:
        doc.status = ParseStatus.NO_TEXT_LAYER
        doc.summary = f"{doc.filename}: the file is empty."
        return doc

    body, was_cut = clip(text, MAX_DOCUMENT_CHARS)
    doc.segments = [Segment(ref=SegmentRef(RefKind.WHOLE), text=body, truncated=was_cut)]
    doc.metadata |= {"characters": len(text), "lines": text.count("\n") + 1}
    doc.status = ParseStatus.PARTIAL if was_cut else ParseStatus.PARSED
    doc.summary = f"{doc.filename}: text file, {len(text)} characters."
    if was_cut:
        doc.warnings.append("The file was longer than the context budget and was truncated.")
    return doc


# ---------------------------------------------------------------------------
# Image
# ---------------------------------------------------------------------------

MAX_IMAGE_EDGE: Final = 1536


def parse_image(doc: ParsedDocument, data: bytes) -> ParsedDocument:
    from PIL import Image, ImageOps, UnidentifiedImageError

    try:
        opened = Image.open(io.BytesIO(data))
        image: Any = ImageOps.exif_transpose(opened) or opened
        original = (image.width, image.height)
        image.thumbnail((MAX_IMAGE_EDGE, MAX_IMAGE_EDGE), Image.Resampling.LANCZOS)
    except (UnidentifiedImageError, OSError) as exc:
        raise ParseFailure("This image could not be read.") from exc

    output = io.BytesIO()
    transparent = image.mode in {"RGBA", "LA"} or (
        image.mode == "P" and "transparency" in image.info
    )
    if transparent:
        image.convert("RGBA").save(output, format="PNG", optimize=True)
        mime = "image/png"
    else:
        image.convert("RGB").save(output, format="JPEG", quality=88, optimize=True)
        mime = "image/jpeg"

    processed = output.getvalue()
    doc.mime_type = mime
    doc.image_data_url = f"data:{mime};base64,{base64.b64encode(processed).decode()}"
    doc.metadata |= {
        "width": image.width,
        "height": image.height,
        "original_width": original[0],
        "original_height": original[1],
        "encoded_bytes": len(processed),
    }
    doc.status = ParseStatus.PARSED
    doc.summary = f"{doc.filename}: image, {image.width}\u00d7{image.height}."
    return doc


PARSERS: Final[dict[DocumentKind, Any]] = {
    DocumentKind.PDF: parse_pdf,
    DocumentKind.DOCX: parse_docx,
    DocumentKind.PPTX: parse_pptx,
    DocumentKind.XLSX: parse_xlsx,
    DocumentKind.CSV: parse_csv,
    DocumentKind.TEXT: parse_text,
    DocumentKind.IMAGE: parse_image,
}
