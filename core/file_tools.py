from __future__ import annotations

import base64
import csv
import io
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import pandas as pd
from docx import Document
from openpyxl import load_workbook
from PIL import Image, ImageOps, UnidentifiedImageError
from pptx import Presentation
from pypdf import PdfReader


MAX_IMAGE_EDGE = 2048
MAX_FILE_TEXT_CHARS = 45_000
MAX_TOTAL_CONTEXT_CHARS = 70_000
IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png", ".webp"}
TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".json",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".html",
    ".css",
    ".xml",
    ".yaml",
    ".yml",
    ".srt",
    ".vtt",
    ".log",
}


@dataclass(slots=True)
class ParsedUpload:
    name: str
    kind: str
    text: str = ""
    summary: str = ""
    image_data_url: str = ""
    image_base64: str = ""
    mime_type: str = ""
    size_bytes: int = 0


def _decode_text(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _parse_image(name: str, raw: bytes) -> ParsedUpload:
    try:
        image = Image.open(io.BytesIO(raw))
        image = ImageOps.exif_transpose(image)
        image.thumbnail((MAX_IMAGE_EDGE, MAX_IMAGE_EDGE), Image.Resampling.LANCZOS)
    except (UnidentifiedImageError, OSError) as exc:
        raise ValueError(f"{name} is not a readable image.") from exc

    output = io.BytesIO()
    has_transparency = image.mode in {"RGBA", "LA"} or (
        image.mode == "P" and "transparency" in image.info
    )
    if has_transparency:
        image = image.convert("RGBA")
        image.save(output, format="PNG", optimize=True)
        mime_type = "image/png"
    else:
        image = image.convert("RGB")
        image.save(output, format="JPEG", quality=88, optimize=True)
        mime_type = "image/jpeg"

    processed = output.getvalue()
    encoded = base64.b64encode(processed).decode("utf-8")
    return ParsedUpload(
        name=name,
        kind="image",
        summary=f"Image: {name} ({image.width}×{image.height})",
        image_data_url=f"data:{mime_type};base64,{encoded}",
        image_base64=encoded,
        mime_type=mime_type,
        size_bytes=len(processed),
    )


def _parse_pdf(name: str, raw: bytes) -> ParsedUpload:
    reader = PdfReader(io.BytesIO(raw))
    parts: list[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(f"--- Page {page_number} ---\n{text}")
        if sum(len(part) for part in parts) >= MAX_FILE_TEXT_CHARS:
            break
    text = "\n\n".join(parts)[:MAX_FILE_TEXT_CHARS]
    return ParsedUpload(
        name=name,
        kind="pdf",
        text=text,
        summary=f"PDF: {name} · {len(reader.pages)} page(s) · {len(text):,} extracted characters",
        size_bytes=len(raw),
    )


def _parse_docx(name: str, raw: bytes) -> ParsedUpload:
    document = Document(io.BytesIO(raw))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        value = paragraph.text.strip()
        if value:
            parts.append(value)
    for table_index, table in enumerate(document.tables, start=1):
        parts.append(f"\n[TABLE {table_index}]")
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    text = "\n".join(parts)[:MAX_FILE_TEXT_CHARS]
    return ParsedUpload(
        name=name,
        kind="docx",
        text=text,
        summary=f"Word document: {name} · {len(text):,} extracted characters",
        size_bytes=len(raw),
    )


def _parse_pptx(name: str, raw: bytes) -> ParsedUpload:
    presentation = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"--- Slide {slide_number} ---\n" + "\n".join(slide_parts))
    text = "\n\n".join(parts)[:MAX_FILE_TEXT_CHARS]
    return ParsedUpload(
        name=name,
        kind="pptx",
        text=text,
        summary=f"PowerPoint: {name} · {len(presentation.slides)} slide(s) · {len(text):,} extracted characters",
        size_bytes=len(raw),
    )


def _parse_csv(name: str, raw: bytes) -> ParsedUpload:
    text_raw = _decode_text(raw)
    dataframe = pd.read_csv(io.StringIO(text_raw), nrows=500)
    preview = dataframe.head(80).to_csv(index=False)
    text = (
        f"Columns: {', '.join(str(column) for column in dataframe.columns)}\n"
        f"Rows loaded: {len(dataframe)}\n\n"
        f"CSV PREVIEW:\n{preview}"
    )[:MAX_FILE_TEXT_CHARS]
    return ParsedUpload(
        name=name,
        kind="csv",
        text=text,
        summary=f"CSV: {name} · {len(dataframe)} loaded row(s) × {len(dataframe.columns)} column(s)",
        size_bytes=len(raw),
    )


def _parse_xlsx(name: str, raw: bytes) -> ParsedUpload:
    workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts: list[str] = []
    sheet_summaries: list[str] = []

    for worksheet in workbook.worksheets[:12]:
        rows: list[list[str]] = []
        for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
            values = ["" if value is None else str(value) for value in row[:30]]
            rows.append(values)
            if row_index >= 150:
                break

        sheet_summaries.append(
            f"{worksheet.title}: {worksheet.max_row} row(s) × {worksheet.max_column} column(s)"
        )
        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerows(rows)
        parts.append(f"--- Sheet: {worksheet.title} ---\n{output.getvalue()}")

    text = "\n\n".join(parts)[:MAX_FILE_TEXT_CHARS]
    return ParsedUpload(
        name=name,
        kind="xlsx",
        text=text,
        summary=f"Excel workbook: {name} · " + "; ".join(sheet_summaries),
        size_bytes=len(raw),
    )


def _parse_json(name: str, raw: bytes) -> ParsedUpload:
    decoded = _decode_text(raw)
    try:
        value = json.loads(decoded)
        text = json.dumps(value, ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        text = decoded
    text = text[:MAX_FILE_TEXT_CHARS]
    return ParsedUpload(
        name=name,
        kind="json",
        text=text,
        summary=f"JSON/text file: {name} · {len(text):,} characters",
        size_bytes=len(raw),
    )


def parse_upload(uploaded_file: Any) -> ParsedUpload:
    name = str(uploaded_file.name)
    suffix = Path(name).suffix.lower()
    raw = uploaded_file.getvalue()

    if suffix in IMAGE_SUFFIXES:
        return _parse_image(name, raw)
    if suffix == ".pdf":
        return _parse_pdf(name, raw)
    if suffix == ".docx":
        return _parse_docx(name, raw)
    if suffix == ".pptx":
        return _parse_pptx(name, raw)
    if suffix == ".csv":
        return _parse_csv(name, raw)
    if suffix in {".xlsx", ".xlsm"}:
        return _parse_xlsx(name, raw)
    if suffix == ".json":
        return _parse_json(name, raw)
    if suffix in TEXT_SUFFIXES:
        text = _decode_text(raw)[:MAX_FILE_TEXT_CHARS]
        return ParsedUpload(
            name=name,
            kind="text",
            text=text,
            summary=f"Text/code file: {name} · {len(text):,} characters",
            size_bytes=len(raw),
        )

    raise ValueError(f"Unsupported file type: {suffix or 'unknown'}")


def build_attachment_context(uploads: list[ParsedUpload]) -> str:
    sections: list[str] = []
    used = 0
    for index, upload in enumerate(uploads, start=1):
        if not upload.text:
            continue
        remaining = MAX_TOTAL_CONTEXT_CHARS - used
        if remaining <= 0:
            break
        text = upload.text[:remaining]
        sections.append(
            f"[FILE {index}] {upload.name}\nTYPE: {upload.kind}\nCONTENT:\n{text}"
        )
        used += len(text)

    if not sections:
        return ""
    return (
        "UPLOADED FILE CONTENT\n"
        "Use this content to answer the user's request. Mention when extraction may "
        "have omitted visual formatting, charts, scanned text, or unsupported elements.\n\n"
        + "\n\n---\n\n".join(sections)
    )
